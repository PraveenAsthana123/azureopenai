"""
Document Parser Module
Supports PDF, DOCX, PPTX, HTML, Markdown extraction with structure preservation.
Uses Azure Document Intelligence for complex documents.
"""

import io
import logging
import re
from abc import ABC, abstractmethod
from dataclasses import dataclass, field
from typing import BinaryIO

from azure.ai.formrecognizer.aio import DocumentAnalysisClient
from azure.identity.aio import DefaultAzureCredential
from bs4 import BeautifulSoup
import markdown
from docx import Document as DocxDocument
from pptx import Presentation
from pypdf import PdfReader

logger = logging.getLogger(__name__)


@dataclass
class ParsedSection:
    """Represents a section of parsed content."""
    heading: str
    level: int
    content: str
    page_number: int | None = None
    bounding_box: list[float] | None = None


@dataclass
class ParsedDocument:
    """Represents a fully parsed document."""
    title: str
    content: str
    sections: list[ParsedSection] = field(default_factory=list)
    pages: list[str] = field(default_factory=list)
    tables: list[dict] = field(default_factory=list)
    images: list[dict] = field(default_factory=list)
    metadata: dict = field(default_factory=dict)
    language: str = "en"
    total_pages: int = 0


class BaseParser(ABC):
    """Base class for document parsers."""

    @abstractmethod
    async def parse(self, content: bytes, filename: str) -> ParsedDocument:
        """Parse document content."""
        pass

    def _clean_text(self, text: str) -> str:
        """Clean extracted text."""
        # Remove excessive whitespace
        text = re.sub(r"\s+", " ", text)
        # Remove control characters
        text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x9f]", "", text)
        return text.strip()


class PDFParser(BaseParser):
    """Parser for PDF documents using pypdf."""

    async def parse(self, content: bytes, filename: str) -> ParsedDocument:
        """Parse PDF content."""
        pdf_file = io.BytesIO(content)
        reader = PdfReader(pdf_file)

        pages = []
        full_text = []
        sections = []
        current_heading = ""

        for page_num, page in enumerate(reader.pages, 1):
            text = page.extract_text() or ""
            pages.append(text)
            full_text.append(text)

            # Try to detect headings (lines in all caps or bold)
            lines = text.split("\n")
            for line in lines:
                line = line.strip()
                if line and (line.isupper() or len(line) < 100):
                    # Potential heading
                    if line.isupper() and len(line) > 3:
                        current_heading = line.title()
                        sections.append(ParsedSection(
                            heading=current_heading,
                            level=1,
                            content="",
                            page_number=page_num,
                        ))

        # Combine all text
        combined_text = "\n\n".join(full_text)

        # Extract metadata
        metadata = {}
        if reader.metadata:
            metadata = {
                "author": reader.metadata.get("/Author", ""),
                "title": reader.metadata.get("/Title", ""),
                "subject": reader.metadata.get("/Subject", ""),
                "creator": reader.metadata.get("/Creator", ""),
                "creation_date": str(reader.metadata.get("/CreationDate", "")),
            }

        return ParsedDocument(
            title=metadata.get("title", filename),
            content=self._clean_text(combined_text),
            sections=sections,
            pages=pages,
            metadata=metadata,
            total_pages=len(reader.pages),
        )


class DocxParser(BaseParser):
    """Parser for Microsoft Word documents."""

    async def parse(self, content: bytes, filename: str) -> ParsedDocument:
        """Parse DOCX content."""
        doc_file = io.BytesIO(content)
        doc = DocxDocument(doc_file)

        sections = []
        full_text = []
        tables = []
        current_section_content = []
        current_heading = ""
        current_level = 0

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            # Check if it's a heading
            if para.style and para.style.name.startswith("Heading"):
                # Save previous section
                if current_heading and current_section_content:
                    sections.append(ParsedSection(
                        heading=current_heading,
                        level=current_level,
                        content="\n".join(current_section_content),
                    ))
                    current_section_content = []

                # Extract heading level
                try:
                    current_level = int(para.style.name.replace("Heading ", ""))
                except ValueError:
                    current_level = 1

                current_heading = text
            else:
                current_section_content.append(text)
                full_text.append(text)

        # Save last section
        if current_heading and current_section_content:
            sections.append(ParsedSection(
                heading=current_heading,
                level=current_level,
                content="\n".join(current_section_content),
            ))

        # Extract tables
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                table_data.append(row_data)
            if table_data:
                tables.append({
                    "headers": table_data[0] if table_data else [],
                    "rows": table_data[1:] if len(table_data) > 1 else [],
                })

        # Extract metadata from core properties
        metadata = {}
        if doc.core_properties:
            metadata = {
                "author": doc.core_properties.author or "",
                "title": doc.core_properties.title or "",
                "subject": doc.core_properties.subject or "",
                "created": str(doc.core_properties.created or ""),
                "modified": str(doc.core_properties.modified or ""),
            }

        return ParsedDocument(
            title=metadata.get("title", filename),
            content=self._clean_text("\n\n".join(full_text)),
            sections=sections,
            tables=tables,
            metadata=metadata,
        )


class PptxParser(BaseParser):
    """Parser for PowerPoint presentations."""

    async def parse(self, content: bytes, filename: str) -> ParsedDocument:
        """Parse PPTX content."""
        ppt_file = io.BytesIO(content)
        prs = Presentation(ppt_file)

        sections = []
        full_text = []
        pages = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            slide_title = ""

            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        # First text box is usually the title
                        if not slide_title and shape.has_text_frame:
                            slide_title = text
                        slide_text.append(text)

                # Extract table content
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells)
                        slide_text.append(row_text)

            page_content = "\n".join(slide_text)
            pages.append(page_content)
            full_text.append(page_content)

            if slide_title:
                sections.append(ParsedSection(
                    heading=slide_title,
                    level=1,
                    content="\n".join(slide_text[1:]) if len(slide_text) > 1 else "",
                    page_number=slide_num,
                ))

        return ParsedDocument(
            title=filename,
            content=self._clean_text("\n\n".join(full_text)),
            sections=sections,
            pages=pages,
            total_pages=len(prs.slides),
        )


class HtmlParser(BaseParser):
    """Parser for HTML documents."""

    async def parse(self, content: bytes, filename: str) -> ParsedDocument:
        """Parse HTML content."""
        html_text = content.decode("utf-8", errors="ignore")
        soup = BeautifulSoup(html_text, "html.parser")

        # Remove unwanted elements
        for element in soup(["script", "style", "nav", "footer", "header", "aside"]):
            element.decompose()

        sections = []
        full_text = []

        # Extract title
        title = ""
        title_tag = soup.find("title")
        if title_tag:
            title = title_tag.get_text(strip=True)

        # Extract headings and their content
        for heading in soup.find_all(["h1", "h2", "h3", "h4", "h5", "h6"]):
            level = int(heading.name[1])
            heading_text = heading.get_text(strip=True)

            # Get content until next heading
            content_parts = []
            for sibling in heading.find_next_siblings():
                if sibling.name in ["h1", "h2", "h3", "h4", "h5", "h6"]:
                    break
                text = sibling.get_text(strip=True)
                if text:
                    content_parts.append(text)

            sections.append(ParsedSection(
                heading=heading_text,
                level=level,
                content="\n".join(content_parts),
            ))

        # Get full text
        full_text = soup.get_text(separator="\n", strip=True)

        # Extract metadata
        metadata = {}
        for meta in soup.find_all("meta"):
            name = meta.get("name", meta.get("property", ""))
            content_val = meta.get("content", "")
            if name and content_val:
                metadata[name] = content_val

        return ParsedDocument(
            title=title or filename,
            content=self._clean_text(full_text),
            sections=sections,
            metadata=metadata,
        )


class MarkdownParser(BaseParser):
    """Parser for Markdown documents."""

    HEADING_PATTERN = re.compile(r"^(#{1,6})\s+(.+)$", re.MULTILINE)

    async def parse(self, content: bytes, filename: str) -> ParsedDocument:
        """Parse Markdown content."""
        md_text = content.decode("utf-8", errors="ignore")

        sections = []
        current_heading = ""
        current_level = 0
        current_content = []

        lines = md_text.split("\n")

        for line in lines:
            heading_match = self.HEADING_PATTERN.match(line)

            if heading_match:
                # Save previous section
                if current_heading:
                    sections.append(ParsedSection(
                        heading=current_heading,
                        level=current_level,
                        content="\n".join(current_content),
                    ))
                    current_content = []

                current_level = len(heading_match.group(1))
                current_heading = heading_match.group(2).strip()
            else:
                current_content.append(line)

        # Save last section
        if current_heading:
            sections.append(ParsedSection(
                heading=current_heading,
                level=current_level,
                content="\n".join(current_content),
            ))

        # Convert to HTML then extract plain text
        html = markdown.markdown(md_text)
        soup = BeautifulSoup(html, "html.parser")
        plain_text = soup.get_text(separator="\n", strip=True)

        # Extract title from first H1
        title = filename
        if sections and sections[0].level == 1:
            title = sections[0].heading

        return ParsedDocument(
            title=title,
            content=self._clean_text(plain_text),
            sections=sections,
        )


class AzureDocumentIntelligenceParser(BaseParser):
    """
    Parser using Azure Document Intelligence for complex documents.
    Provides better accuracy for scanned PDFs, forms, and complex layouts.
    """

    def __init__(self, endpoint: str):
        self.endpoint = endpoint
        self.credential = None
        self.client = None

    async def connect(self):
        """Initialize Document Intelligence client."""
        self.credential = DefaultAzureCredential()
        self.client = DocumentAnalysisClient(
            endpoint=self.endpoint,
            credential=self.credential,
        )

    async def close(self):
        """Close connections."""
        if self.client:
            await self.client.close()
        if self.credential:
            await self.credential.close()

    async def parse(self, content: bytes, filename: str) -> ParsedDocument:
        """Parse document using Azure Document Intelligence."""
        if not self.client:
            await self.connect()

        # Use prebuilt-layout model for best structure extraction
        poller = await self.client.begin_analyze_document(
            "prebuilt-layout",
            document=content,
        )
        result = await poller.result()

        sections = []
        pages = []
        tables = []
        full_text = []

        # Process pages
        for page in result.pages:
            page_text = []

            # Extract lines
            for line in page.lines:
                page_text.append(line.content)
                full_text.append(line.content)

            pages.append("\n".join(page_text))

        # Extract paragraphs as sections
        if result.paragraphs:
            for para in result.paragraphs:
                # Check if it's a heading based on role
                if para.role and "heading" in para.role.lower():
                    level = 1 if "title" in para.role.lower() else 2
                    sections.append(ParsedSection(
                        heading=para.content,
                        level=level,
                        content="",
                        page_number=para.bounding_regions[0].page_number if para.bounding_regions else None,
                    ))

        # Extract tables
        for table in result.tables:
            table_data = {
                "row_count": table.row_count,
                "column_count": table.column_count,
                "cells": [],
            }
            for cell in table.cells:
                table_data["cells"].append({
                    "row": cell.row_index,
                    "column": cell.column_index,
                    "content": cell.content,
                    "is_header": cell.kind == "columnHeader",
                })
            tables.append(table_data)

        return ParsedDocument(
            title=filename,
            content=self._clean_text("\n\n".join(full_text)),
            sections=sections,
            pages=pages,
            tables=tables,
            total_pages=len(result.pages),
            language=result.languages[0] if result.languages else "en",
        )


class DocumentParserFactory:
    """Factory for creating document parsers."""

    def __init__(self, doc_intelligence_endpoint: str | None = None):
        self.doc_intelligence_endpoint = doc_intelligence_endpoint
        self._parsers = {
            ".pdf": PDFParser(),
            ".docx": DocxParser(),
            ".doc": DocxParser(),
            ".pptx": PptxParser(),
            ".ppt": PptxParser(),
            ".html": HtmlParser(),
            ".htm": HtmlParser(),
            ".md": MarkdownParser(),
            ".markdown": MarkdownParser(),
            ".txt": MarkdownParser(),
        }

    def get_parser(self, filename: str, use_doc_intelligence: bool = False) -> BaseParser:
        """Get appropriate parser for file type."""
        ext = "." + filename.split(".")[-1].lower() if "." in filename else ""

        if use_doc_intelligence and self.doc_intelligence_endpoint:
            return AzureDocumentIntelligenceParser(self.doc_intelligence_endpoint)

        parser = self._parsers.get(ext)
        if not parser:
            raise ValueError(f"No parser available for file type: {ext}")

        return parser

    async def parse(
        self,
        content: bytes,
        filename: str,
        use_doc_intelligence: bool = False,
    ) -> ParsedDocument:
        """Parse document using appropriate parser."""
        parser = self.get_parser(filename, use_doc_intelligence)

        if isinstance(parser, AzureDocumentIntelligenceParser):
            await parser.connect()
            try:
                return await parser.parse(content, filename)
            finally:
                await parser.close()

        return await parser.parse(content, filename)
